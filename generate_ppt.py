"""
PPT Generator - SPK Hipertensi Project Presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import copy
from lxml import etree

# ── WARNA TEMA ──────────────────────────────────────────────────────────────
COl_BG_DARK   = RGBColor(0x1A, 0x23, 0x5E)   # biru tua
COl_BG_MID    = RGBColor(0x2E, 0x40, 0x8A)   # biru medium
COl_ACCENT    = RGBColor(0x66, 0x7E, 0xEA)   # biru cerah
COl_GOLD      = RGBColor(0xF3, 0x9C, 0x12)   # gold/oren
COl_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
COl_LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xFF)   # biru sangat muda
COl_CARD      = RGBColor(0xE8, 0xEF, 0xFF)   # card background
COl_TEXT_DARK = RGBColor(0x1A, 0x1A, 0x2E)   # teks gelap
COl_TEXT_MID  = RGBColor(0x4A, 0x4A, 0x6A)   # teks sedang
COl_GREEN     = RGBColor(0x27, 0xAE, 0x60)   # hijau
COl_RED       = RGBColor(0xE7, 0x4C, 0x3C)   # merah
COl_BAR_GOOD  = RGBColor(0x2E, 0xCC, 0x71)   # hijau bar
COl_BAR_BAD   = RGBColor(0xE7, 0x4C, 0x3C)   # merah bar


def set_shape_fill(shape, rgb: RGBColor):
    """Isi shape dengan warna solid."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def set_shape_line(shape, rgb: RGBColor, width_pt=0):
    """Hapus border shape."""
    line = shape.line
    if width_pt == 0:
        line.fill.background()
    else:
        line.color.rgb = rgb
        line.width = Pt(width_pt)


def add_textbox(slide, left, top, width, height,
                text, font_size=18, bold=False,
                color=COl_TEXT_DARK, align=PP_ALIGN.LEFT,
                wrap=True, italic=False):
    """Helper membuat textbox."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_rect(slide, left, top, width, height, fill_rgb, line_rgb=None, radius=False):
    """Helper membuat rectangle / rounded rectangle."""
    from pptx.util import Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    if radius:
        shape = slide.shapes.add_shape(
            5,  # MSO_SHAPE_TYPE.ROUNDED_RECTANGLE
            left, top, width, height
        )
    else:
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            left, top, width, height
        )
    set_shape_fill(shape, fill_rgb)
    if line_rgb:
        set_shape_line(shape, line_rgb, 1)
    else:
        set_shape_line(shape, fill_rgb, 0)
    return shape


def add_rect_text(slide, left, top, width, height,
                  text, fill_rgb, text_color=COl_WHITE,
                  font_size=16, bold=True, align=PP_ALIGN.CENTER,
                  v_anchor=None, italic=False):
    """Rectangle dengan teks di dalamnya."""
    rect = add_rect(slide, left, top, width, height, fill_rgb, radius=True)
    tf = rect.text_frame
    tf.word_wrap = True
    if v_anchor:
        tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    run.font.italic = italic
    return rect


def add_slide_header(slide, title, subtitle=None,
                     bg_left=COl_BG_DARK, bg_right=COl_BG_MID):
    """Header standar untuk setiap slide."""
    W = Inches(13.33)
    H = Inches(1.1)

    # Background gradient simulation (2 rectangles)
    add_rect(slide, Inches(0), Inches(0), W, H, bg_left)
    rect_r = add_rect(slide, Inches(4.5), Inches(0), Inches(8.83), H, bg_right)

    # Accent stripe
    add_rect(slide, Inches(0), Inches(1.05), W, Inches(0.06), COl_GOLD)

    # Title
    add_textbox(slide, Inches(0.4), Inches(0.1), Inches(9), Inches(0.7),
                title, font_size=28, bold=True, color=COl_WHITE)

    # Subtitle
    if subtitle:
        add_textbox(slide, Inches(0.4), Inches(0.75), Inches(9), Inches(0.4),
                    subtitle, font_size=14, color=RGBColor(0xCC, 0xD6, 0xFF))

    # Logo badge
    add_rect_text(slide, Inches(11.8), Inches(0.15), Inches(1.3), Inches(0.7),
                  "SPK\nHipertensi", COl_ACCENT, COl_WHITE, font_size=11, bold=True)


def pptx_table(slide, data, col_widths, left, top,
                header_fill=COl_BG_DARK, row_fill=COl_LIGHT_BG,
                alt_fill=COl_WHITE, font_size=11, header_font_size=12):
    """Buat tabel PPTX."""
    rows = len(data)
    cols = len(data[0])
    total_w = sum(col_widths)

    tbl = slide.shapes.add_table(rows, cols, left, top,
                                 total_w, Inches(0.4 * rows)).table

    # Set column widths
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    for ri, row_data in enumerate(data):
        for ci, cell_text in enumerate(row_data):
            cell = tbl.cell(ri, ci)
            cell.text = str(cell_text)

            # Style
            tf = cell.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER

            run = p.runs[0] if p.runs else p.add_run()
            run.text = str(cell_text)
            run.font.size = Pt(font_size if ri > 0 else header_font_size)
            run.font.bold = (ri == 0)
            run.font.color.rgb = COl_WHITE if ri == 0 else COl_TEXT_DARK

            # Fill
            from pptx.oxml.ns import qn
            from lxml import etree
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
            srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
            def rgb_to_hex(c):
                return f'{c[0]:02X}{c[1]:02X}{c[2]:02X}'
            if ri == 0:
                srgbClr.set('val', rgb_to_hex(header_fill))
            elif ri % 2 == 0:
                srgbClr.set('val', rgb_to_hex(alt_fill))
            else:
                srgbClr.set('val', rgb_to_hex(row_fill))

    return tbl


# ════════════════════════════════════════════════════════════════════════════
#  BUILD PRESENTATION
# ════════════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # blank

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 - COVER
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
W = prs.slide_width
H = prs.slide_height

# Background gradient
add_rect(sl, 0, 0, W, H, COl_BG_DARK)
add_rect(sl, Inches(4.5), 0, Inches(8.83), H, COl_BG_MID)

# Accent lines
add_rect(sl, 0, Inches(2.0), W, Inches(0.08), COl_GOLD)
add_rect(sl, 0, Inches(4.8), W, Inches(0.08), COl_GOLD)

# Main title
add_textbox(sl, Inches(0.8), Inches(2.2), Inches(11), Inches(1.0),
            "SISTEM PENDUKUNG KEPUTUSAN",
            font_size=40, bold=True, color=COl_GOLD, align=PP_ALIGN.CENTER)

add_textbox(sl, Inches(0.8), Inches(3.15), Inches(11), Inches(1.2),
            "DETEKSI HIPERTENSI MENGGUNAKAN\nDECISION TREE DAN PROBABILITY CALIBRATION",
            font_size=28, bold=True, color=COl_WHITE, align=PP_ALIGN.CENTER)

add_textbox(sl, Inches(0.8), Inches(4.9), Inches(11), Inches(0.5),
            "Daniko Sutopo",
            font_size=20, color=COl_WHITE, align=PP_ALIGN.CENTER)

add_textbox(sl, Inches(0.8), Inches(5.4), Inches(11), Inches(0.4),
            "Teknik Informatika - Sistem Pendukung Keputusan",
            font_size=14, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.CENTER)

# Icon placeholder - hospital cross
icon_bg = add_rect_text(sl, Inches(10.9), Inches(2.5), Inches(1.9), Inches(1.9),
                         "🏥", COl_ACCENT, COl_WHITE, font_size=48, bold=False)
icon_bg.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 - DAFTAR ISI
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
add_slide_header(sl, "Daftar Isi Presentasi",
                 "Struktur penyampaian materi proyek akhir")

items = [
    ("01", "Dataset",             "Sumber data, fitur, dan ukuran dataset"),
    ("02", "Goals & Kriteria",    "Tujuan sistem, kriteria keputusan, dan alternatif"),
    ("03", "Metode",             "Decision Tree C4.5 + Probability Calibration"),
    ("04", "Pengolahan Data",     "Pra-pemrosesan, encoding, training, evaluasi"),
    ("05", "Demo Program",       "Tampilan aplikasi dan alur prediksi"),
]

y_start = Inches(1.35)
for i, (num, title, desc) in enumerate(items):
    y = y_start + i * Inches(1.05)
    # Number box
    add_rect_text(sl, Inches(0.6), y, Inches(0.9), Inches(0.8),
                  num, COl_ACCENT, COl_WHITE, font_size=20, bold=True)
    # Title
    add_textbox(sl, Inches(1.7), y + Inches(0.05), Inches(5), Inches(0.45),
                title, font_size=20, bold=True, color=COl_TEXT_DARK)
    # Desc
    add_textbox(sl, Inches(1.7), y + Inches(0.45), Inches(8), Inches(0.35),
                desc, font_size=13, color=COl_TEXT_MID)

    # Connector line
    if i < len(items) - 1:
        add_rect(sl, Inches(1.05), y + Inches(0.82),
                 Inches(0.02), Inches(0.22), COl_ACCENT)

# Side accent
add_rect(sl, Inches(11.8), Inches(1.3), Inches(1.3), Inches(5.8), COl_LIGHT_BG)
add_rect_text(sl, Inches(11.85), Inches(3.3), Inches(1.2), Inches(1.0),
              "5\nPOIN", COl_BG_DARK, COl_GOLD, font_size=22, bold=True)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 - DATASET
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
add_slide_header(sl, "01  Dataset", "Sumber dan spesifikasi data yang digunakan")

# Info cards row
card_data = [
    ("📁 Sumber Data",  "D:\\SPK HYPERTENSION\\hypertension_dataset.xlsx"),
    ("📊 Total Sample",  "1.985 record"),
    ("🔢 Fitur",       "9 fitur (input) + 1 target"),
    ("🎯 Target",      "Has_Hypertension (0 = Tidak, 1 = Ya)"),
]
cx = Inches(0.4)
for label, val in card_data:
    add_rect_text(sl, cx, Inches(1.3), Inches(3.0), Inches(1.2),
                  f"{label}\n{val}", COl_CARD, COl_TEXT_DARK,
                  font_size=13, bold=False)
    cx += Inches(3.2)

# Feature table
feat_header = ["No", "Fitur", "Tipe Data", "Nilai Original", "Encoded"]
feat_rows = [
    ["1",  "Age",              "Numerik",  "1 - 120 tahun",           "digunakan langsung"],
    ["2",  "Salt_Intake",       "Numerik",  "0 - 25 gram/hari",         "digunakan langsung"],
    ["3",  "Stress_Score",      "Numerik",  "0 - 10",                   "digunakan langsung"],
    ["4",  "BP_History",       "Kategorik","Normal / Prehyper / Hyper","1 / 2 / 3"],
    ["5",  "Sleep_Duration",    "Numerik",  "0 - 24 jam/hari",          "digunakan langsung"],
    ["6",  "BMI",               "Numerik",  "10 - 80",                  "digunakan langsung"],
    ["7",  "Family_History",    "Kategorik","No / Yes",                 "0 / 1"],
    ["8",  "Smoking_Status",    "Kategorik","Non-Smoker / Smoker",      "0 / 1"],
    ["9",  "Exercise_Level",     "Kategorik","Low / Moderate / High",   "1 / 2 / 3"],
]
feat_table = [feat_header] + feat_rows

pptx_table(sl, feat_table,
            [Inches(0.5), Inches(2.3), Inches(1.6), Inches(3.0), Inches(2.8)],
            Inches(0.4), Inches(2.7),
            header_fill=COl_BG_DARK, font_size=11, header_font_size=12)

# Note
add_textbox(sl, Inches(0.4), Inches(7.0), Inches(12), Inches(0.4),
            "📌 Split data: 80% training (1.588) | 20% testing (397) - stratified sampling",
            font_size=12, color=COl_TEXT_MID, italic=True)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 - GOALS & KRITERIA
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
add_slide_header(sl, "02  Goals & Kriteria", "Tujuan, variabel keputusan, dan alternatif")

# LEFT - Goals
add_rect_text(sl, Inches(0.4), Inches(1.3), Inches(5.8), Inches(0.6),
              "🎯  Tujuan Sistem", COl_BG_DARK, COl_WHITE, font_size=18, bold=True)

goal_items = [
    "Membantu tenaga medis mendeteksi risiko hipertensi secara otomatis",
    "Memberikan hasil prediksi disertai tingkat kepercayaan (confidence)",
    "Menjelaskan alur keputusan secara visual dan mudah dipahami",
    "Meningkatkan akurasi diagnosis dibandingkan metode manual",
]
y = Inches(2.0)
for g in goal_items:
    add_rect(sl, Inches(0.4), y, Inches(0.08), Inches(0.5), COl_GOLD)
    add_textbox(sl, Inches(0.65), y, Inches(5.5), Inches(0.6),
                g, font_size=14, color=COl_TEXT_DARK)
    y += Inches(0.75)

# RIGHT - Criteria & Alternatives
add_rect_text(sl, Inches(6.6), Inches(1.3), Inches(6.3), Inches(0.6),
              "⚖️  Kriteria & Alternatif", COl_BG_MID, COl_WHITE, font_size=18, bold=True)

criteria_data = [
    ("Kriteria",          "Bobot Prioritas"),
    ("Riwayat Tekanan Darah (BP_History)",  "28.04% - Paling Penting"),
    ("Usia (Age)",                     "13.70%"),
    ("Riwayat Keluarga (Family_History)",   "11.76%"),
    ("Tingkat Stres (Stress_Score)",        "10.03%"),
    ("BMI",                            "09.62%"),
    ("Status Merokok (Smoking_Status)",    "09.06%"),
    ("Asupan Garam (Salt_Intake)",          "08.98%"),
    ("Durasi Tidur (Sleep_Duration)",       "08.82%"),
    ("Aktivitas Fisik (Exercise_Level)",    "00.00% - Diabaikan model"),
]
pptx_table(sl, criteria_data,
            [Inches(4.0), Inches(2.1)],
            Inches(6.6), Inches(2.0),
            header_fill=COl_BG_MID, font_size=11, header_font_size=12)

# Alternatives note
add_rect_text(sl, Inches(0.4), Inches(5.1), Inches(6.0), Inches(1.6),
              "💡 Alternatif Keputusan:\n\n"
              "✅  Tidak Hipertensi - gaya hidup normal\n"
              "⚠️  Hipertensi - konsultasi dokter diperlukan",
              COl_LIGHT_BG, COl_TEXT_DARK, font_size=14, bold=False)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 - METODE (Decision Tree)
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
add_slide_header(sl, "03  Metode - Decision Tree", "Algoritma utama yang digunakan")

# Left column - explanation
add_textbox(sl, Inches(0.4), Inches(1.3), Inches(5.5), Inches(0.5),
            "Apa itu Decision Tree?", font_size=20, bold=True, color=COl_BG_DARK)

dt_desc = (
    "Decision Tree (Pohon Keputusan) adalah algoritma machine learning "
    "yang bekerja seperti flowchart: setiap node mewakili sebuah pertanyaan "
    "terhadap satu fitur data, dan setiap cabang menunjukkan jawaban ya/tidak. "
    "Proses ini berulang sampai mencapai daun (leaf) yang berisi hasil akhir."
)
add_textbox(sl, Inches(0.4), Inches(1.85), Inches(5.8), Inches(1.2),
            dt_desc, font_size=13, color=COl_TEXT_MID)

# How it works cards
steps = [
    ("1", "Root Node",       "Mulai dari akar, pilih fitur terbaik\nuntuk memisahkan data"),
    ("2", "Split",           "Bagi data menjadi dua kelompok\nberdasarkan threshold"),
    ("3", "Recursive",       "Ulangi langkah di setiap node\nhingga kriteria berhenti terpenuhi"),
    ("4", "Leaf Node",       "Node akhir berisi kelas prediksi\n(Tidak / Ya Hipertensi)"),
]
cx = Inches(0.4)
for num, title, desc in steps:
    add_rect_text(sl, cx, Inches(3.3), Inches(3.0), Inches(1.7),
                  f"{num}\n{title}\n\n{desc}",
                  COl_CARD, COl_TEXT_DARK, font_size=13, bold=False)
    cx += Inches(3.3)

# Right - Hyperparameters & Metrics
add_rect_text(sl, Inches(6.5), Inches(1.3), Inches(6.5), Inches(0.55),
              "⚙️  Best Hyperparameters (GridSearchCV)", COl_BG_DARK, COl_WHITE,
              font_size=16, bold=True)

hp_data = [
    ("Parameter",           "Nilai",         "Keterangan"),
    ("criterion",          "gini",           "Ukuran kemurnian node"),
    ("max_depth",           "8",              "Kedalaman maksimum pohon"),
    ("min_samples_split",   "5",              "Min sample untuk split"),
    ("min_samples_leaf",    "1",              "Min sample di leaf"),
    ("cv_folds",            "5",              "Cross-validation"),
]
pptx_table(sl, hp_data,
            [Inches(2.4), Inches(1.4), Inches(2.5)],
            Inches(6.5), Inches(1.95),
            header_fill=COl_BG_MID, font_size=12, header_font_size=13)

# Gini formula visual
add_rect_text(sl, Inches(6.5), Inches(4.2), Inches(6.5), Inches(1.5),
              "📐  Gini Impurity (Kriteria Split):\n\n"
              "  Gini = 1 − Σ (pᵢ)²\n\n"
              "  Semakin kecil nilai Gini → node semakin murni → split semakin baik",
              COl_LIGHT_BG, COl_TEXT_DARK, font_size=13, bold=False)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 - METODE (Probability Calibration)
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
add_slide_header(sl, "03  Metode - Probability Calibration",
                 "Isotonic Regression untuk tingkat kepercayaan yang realistis")

# Problem statement
add_rect_text(sl, Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.7),
              "⚠️  Masalah: Decision Tree murni menghasilkan confidence 100% "
              "untuk hampir semua prediksi - tidak realistis!",
              RGBColor(0xFF, 0xF3, 0xE0), COl_TEXT_DARK, font_size=15, bold=False)

# Solution
add_rect_text(sl, Inches(0.4), Inches(2.1), Inches(12.5), Inches(0.65),
              "✅  Solusi: Probability Calibration dengan Isotonic Regression "
              "(CalibratedClassifierCV, method='isotonic', cv=5)",
              COl_LIGHT_BG, COl_TEXT_DARK, font_size=15, bold=False)

# Two column explanation
# Left
add_rect_text(sl, Inches(0.4), Inches(2.95), Inches(5.8), Inches(0.55),
              "📊 Sebelum Kalibrasi (Tanpa Kalibrasi)", COl_RED, COl_WHITE,
              font_size=14, bold=True)

raw_prob_data = [
    ("Confidence",  "Persentase"),
    ("100%",  "Mayoritas prediksi"),
    ("~99%",  "Hampir semua leaf node"),
]
pptx_table(sl, raw_prob_data,
            [Inches(2.4), Inches(3.2)],
            Inches(0.4), Inches(3.6),
            header_fill=COl_RED, font_size=12, header_font_size=13)

# Right
add_rect_text(sl, Inches(6.7), Inches(2.95), Inches(6.2), Inches(0.55),
              "📊 Setelah Kalibrasi Isotonic (Realistis)", COl_GREEN, COl_WHITE,
              font_size=14, bold=True)

calib_data = [
    ("Confidence",  "Jumlah Sample"),
    ("50.9% - 60%",  "sedikit (uncertain)"),
    ("60% - 80%",    "beberapa"),
    ("80% - 98.1%",  "mayoritas"),
]
pptx_table(sl, calib_data,
            [Inches(2.4), Inches(3.6)],
            Inches(6.7), Inches(3.6),
            header_fill=COl_GREEN, font_size=12, header_font_size=13)

# How isotonic works
add_rect_text(sl, Inches(0.4), Inches(5.2), Inches(12.5), Inches(1.8),
              "📖  Cara Kerja Isotonic Regression:\n\n"
              "1. Model Decision Tree asli memprediksi probabilitas mentah dari proporsi kelas di leaf node.\n"
              "2. Isotonic Regression membangun fungsi monoton (tabel koreksi) menggunakan "
              "algoritma PAVA (Pool Adjacent Violators) pada data cross-validation (5-fold).\n"
              "3. Fungsi ini menyesuaikan probabilitas mentah agar mencerminkan "
              "frekuensi nyata kelas di populasi data - hasilnya lebih jujur dan realistis.",
              COl_CARD, COl_TEXT_DARK, font_size=13, bold=False)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 - PENGOLAHAN DATA
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
add_slide_header(sl, "04  Pengolahan Data", "Tahapan lengkap dari data mentah hingga model siap pakai")

pipeline = [
    ("1", "LOAD",          "Load Dataset",       COl_BG_DARK,
     "Baca file Excel dari\nD:\\SPK HYPERTENSION\\\nhypertension_dataset.xlsx"),
    ("2", "PREPROCESS",   "Encoding Label",      COl_BG_MID,
     "Ubah label kategorik\n(Normal→1, No→0, dll)\nmenggunakan label mapping"),
    ("3", "SPLIT",         "Train-Test Split",    COl_ACCENT,
     "80% training (1.588)\n20% testing (397)\nStratified split"),
    ("4", "TUNE",          "GridSearchCV",       COl_ACCENT,
     "Cari hyperparameter\nterbaik dengan\n5-fold CV"),
    ("5", "TRAIN",         "Train Decision Tree", COl_GREEN,
     "Latih model dengan\nhyperparameter terbaik\ncriterion=gini, depth=8"),
    ("6", "CALIBRATE",     "Isotonic Reg.",      COl_GREEN,
     "Kalibrasi probabilitas\ndengan Calibrated\nClassifierCV (cv=5)"),
    ("7", "EVALUATE",      "Evaluasi Model",     COl_GOLD,
     "Accuracy, Precision,\nRecall, F1-Score\nConfusion Matrix"),
]

n = len(pipeline)
W_slide = Inches(13.33)
total_w = W_slide - Inches(0.8)
box_w = total_w / n - Inches(0.08)

y_boxes = Inches(1.4)
y_labels = Inches(2.85)
y_arrows = Inches(2.55)

for i, (num, phase, title, color, desc) in enumerate(pipeline):
    x = Inches(0.4) + i * (box_w + Inches(0.08))
    # Box
    add_rect_text(sl, x, y_boxes, box_w, Inches(1.3),
                  f"{num}\n{phase}\n{title}", color, COl_WHITE,
                  font_size=11, bold=True)
    # Description
    add_textbox(sl, x, y_boxes + Inches(1.35), box_w, Inches(0.9),
                desc, font_size=10, color=COl_TEXT_MID, align=PP_ALIGN.CENTER)
    # Arrow
    if i < n - 1:
        ax = x + box_w + Inches(0.01)
        add_rect(sl, ax, y_arrows, Inches(0.06), Inches(0.35), COl_GOLD)

# Metrics row
add_rect_text(sl, Inches(0.4), Inches(4.3), Inches(12.5), Inches(0.55),
              "📊  Hasil Evaluasi Model", COl_BG_DARK, COl_WHITE,
              font_size=16, bold=True)

metrics = [
    ("Accuracy",  "96.73%", COl_GREEN),
    ("Precision", "97.07%", COl_GREEN),
    ("Recall",    "96.60%", COl_GREEN),
    ("F1-Score",  "96.84%", COl_GREEN),
]
mx = Inches(0.4)
for name, val, col in metrics:
    add_rect_text(sl, mx, Inches(5.0), Inches(3.0), Inches(1.2),
                  f"{name}\n{val}", col, COl_WHITE, font_size=22, bold=True)
    mx += Inches(3.2)

add_textbox(sl, Inches(0.4), Inches(6.3), Inches(12), Inches(0.4),
            "Confidence Distribution: Min 50.9% | Mean 93.3% | Max 98.1% | <90%: 56 sample",
            font_size=13, color=COl_TEXT_MID, italic=True)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 - DEMO PROGRAM
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
add_slide_header(sl, "05  Demo Program", "Tampilan aplikasi Streamlit dan alur prediksi")

# App screenshot placeholder
add_rect_text(sl, Inches(0.4), Inches(1.3), Inches(7.5), Inches(4.5),
              "🖥️\n\nScreenshot Aplikasi\nSPK Hipertensi\n\n(Jalankan: streamlit run spk_app.py)",
              RGBColor(0xE8, 0xEF, 0xFF), COl_TEXT_DARK, font_size=18, bold=False)

# Right - feature list
add_rect_text(sl, Inches(8.2), Inches(1.3), Inches(4.8), Inches(0.55),
              "Fitur Aplikasi", COl_BG_DARK, COl_WHITE, font_size=16, bold=True)

fitur_items = [
    "📋  Form input 9 fitur pasien",
    "🔍  Prediksi real-time + confidence",
    "🌳  Visualisasi alur keputusan (decision path)",
    "📐  Perhitungan probabilitas step-by-step",
    "📊  Feature importance chart",
    "🌲  Pohon keputusan lengkap (full tree)",
    "📄  Laporan training model",
    "📱  Responsive layout (Streamlit)",
]
fy = Inches(1.95)
for item in fitur_items:
    add_textbox(sl, Inches(8.2), fy, Inches(4.8), Inches(0.42),
                item, font_size=13, color=COl_TEXT_DARK)
    fy += Inches(0.44)

# Bottom - prediction flow
add_rect_text(sl, Inches(0.4), Inches(6.0), Inches(12.5), Inches(0.65),
              "🔄  Alur Prediksi: Input → Encode → Decision Tree → Kalibrasi Probabilitas → Hasil + Visualisasi",
              COl_LIGHT_BG, COl_TEXT_DARK, font_size=14, bold=False)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 - KESIMPULAN
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
add_slide_header(sl, "Kesimpulan", "Ringkasan dan hasil akhir proyek")

# Summary cards
summary = [
    ("✅", "Akurasi Tinggi",    "96.73% accuracy dengan Decision Tree + GridSearchCV"),
    ("📊", "Confidence Realistis", "Isotonic Regression membuat probabilitas lebih jujur"),
    ("🌳", "Feature Importance", "BP_History paling berpengaruh (28.04%)"),
    ("👁️", "Visualisasi Jelas",  "Decision path & pohon lengkap dapat dipahami dokter & pasien"),
    ("🔢", "9 Fitur",            "Kombinasi data demografi, gaya hidup, dan riwayat kesehatan"),
    ("📱", "Web App",            "Aplikasi interaktif berbasis Streamlit - mudah digunakan"),
]
sx = Inches(0.4)
sy = Inches(1.5)
for emoji, title, desc in summary:
    add_rect_text(sl, sx, sy, Inches(4.0), Inches(1.5),
                  f"{emoji}\n{title}\n\n{desc}",
                  COl_CARD, COl_TEXT_DARK, font_size=13, bold=False)
    sx += Inches(4.3)
    if sx > Inches(12.5):
        sx = Inches(0.4)
        sy += Inches(1.7)

# Closing
add_rect(sl, 0, Inches(6.0), W, Inches(1.5), COl_BG_DARK)
add_textbox(sl, Inches(0.5), Inches(6.15), Inches(12), Inches(0.6),
            "TERIMA KASIH",
            font_size=32, bold=True, color=COl_GOLD, align=PP_ALIGN.CENTER)
add_textbox(sl, Inches(0.5), Inches(6.75), Inches(12), Inches(0.4),
            "SPK Hipertensi - Decision Tree + Probability Calibration | Daniko Sutopo",
            font_size=14, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────────────────────
output_path = "D:/SPK-HYPERTENSION-APP/SPK_Hipertensi_Presentasi.pptx"
prs.save(output_path)
print(f"OK PPT berhasil disimpan: {output_path}")
print(f"   Total slide: {len(prs.slides)}")
